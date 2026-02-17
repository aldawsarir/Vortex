import os
import re
from flask import Flask, render_template, request, redirect, url_for, session, flash, Response, send_file
from flask_login import LoginManager, login_user, logout_user, login_required, current_user
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash, check_password_hash
from datetime import datetime
from io import BytesIO
import threading
from concurrent.futures import ThreadPoolExecutor

from models import db, User, Upload, QuizResult, DailyChallenge, QuizBattle, BattleParticipant, Review, PuzzleGame, Subject, Chapter
from utils.preprocessing import preprocess_text
from utils.summarizer import summarize_text, summarize_text_with_style, extract_keywords
from utils.quiz import generate_quiz
from utils.visualize import create_mindmap
from utils.visualization import create_bar_chart, create_pie_chart, analyze_quiz_performance
from utils.analytics import calculate_user_stats, generate_performance_report
from utils.ocr import extract_text_from_image
from utils.knowledge_base import get_all_categories, get_topics_by_category, get_topic_content, search_topics
# ================= Configuration =================
ALLOWED_EXTENSIONS = {'pdf', 'docx', 'pptx', 'txt', 'png', 'jpg', 'jpeg', 'bmp', 'tiff', 'gif', 'webp'}
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOAD_FOLDER = os.path.join(BASE_DIR, 'uploads')
STATIC_FOLDER = os.path.join(BASE_DIR, 'static')
TEMPLATES_FOLDER = os.path.join(BASE_DIR, 'templates')

app = Flask(__name__, static_folder=STATIC_FOLDER, template_folder=TEMPLATES_FOLDER)
app.secret_key = "vortex-secret-key-2026-pastel-theme"
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///vortex.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 32 * 1024 * 1024

db.init_app(app)
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login'

limiter = Limiter(
    app=app,
    key_func=get_remote_address,
    default_limits=["200 per day", "50 per hour"],
    storage_uri="memory://",
    headers_enabled=True
)

executor = ThreadPoolExecutor(max_workers=4)

for folder in [UPLOAD_FOLDER, os.path.join(STATIC_FOLDER, 'images')]:
    os.makedirs(folder, exist_ok=True)

with app.app_context():
    db.create_all()

@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_file(filepath, ext):
    try:
        file_text = ""
        
        if ext == 'txt':
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                file_text = f.read(10000)
            print(f"✅ TXT extracted: {len(file_text)} chars")
        
        elif ext == 'pdf':
            import PyPDF2
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                total_pages = len(reader.pages)
                pages_to_read = min(15, total_pages)
                
                text_chunks = []
                for i in range(pages_to_read):
                    try:
                        page_text = reader.pages[i].extract_text() or ""
                        if page_text.strip():
                            text_chunks.append(page_text)
                    except:
                        continue
                
                file_text = "\n".join(text_chunks)[:10000]
            print(f"✅ PDF extracted: {len(file_text)} chars ({pages_to_read}/{total_pages} pages)")
        
        elif ext == 'docx':
            import docx
            doc = docx.Document(filepath)
            
            paragraphs = []
            for i, para in enumerate(doc.paragraphs):
                if i >= 100:
                    break
                text = para.text.strip()
                if text and len(text) > 10:
                    paragraphs.append(text)
            
            file_text = "\n".join(paragraphs)[:10000]
            print(f"✅ DOCX extracted: {len(file_text)} chars ({len(paragraphs)} paragraphs)")
        
        elif ext == 'pptx':
            from pptx import Presentation
            prs = Presentation(filepath)
            total_slides = len(prs.slides)
            slides_to_read = min(10, total_slides)
            extracted_text = []
            
            for slide_num, slide in enumerate(prs.slides[:slides_to_read]):
                if len(extracted_text) >= 50:
                    break
                
                for shape in slide.shapes:
                    if not hasattr(shape, "text"):
                        continue
                    
                    text_content = shape.text.strip()
                    if not text_content or len(text_content) > 300:
                        continue
                    
                    if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                        if len(shape.text_frame.paragraphs) == 1:
                            extracted_text.append(text_content)
                        else:
                            for para in shape.text_frame.paragraphs[:5]:
                                bullet_text = para.text.strip()
                                if bullet_text and len(bullet_text) > 5:
                                    extracted_text.append(bullet_text)
            
            file_text = "\n".join(extracted_text)[:10000]
            print(f"✅ PPTX extracted: {len(file_text)} chars ({slides_to_read}/{total_slides} slides)")
        
        return file_text
    
    except Exception as e:
        print(f"❌ Error extracting from {ext}: {str(e)}")
        return ""

# ================= Authentication Routes =================
@app.route('/register', methods=['GET', 'POST'])
@limiter.limit("5 per minute")
def register():
    if request.method == 'POST':
        username = request.form.get('username', '').strip()
        email = request.form.get('email', '').strip()
        password = request.form.get('password', '')
        
        if not username or not email or not password:
            flash('❌ All fields are required', 'danger')
            return redirect(url_for('register'))
        
        if len(username) < 3 or len(username) > 20:
            flash('❌ Username must be between 3-20 characters', 'danger')
            return redirect(url_for('register'))
        
        if not username.isalnum():
            flash('❌ Username must contain only letters and numbers', 'danger')
            return redirect(url_for('register'))
        
        email_regex = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'
        if not re.match(email_regex, email):
            flash('❌ Please enter a valid email address', 'danger')
            return redirect(url_for('register'))
        
        if len(password) < 8:
            flash('❌ Password must be at least 8 characters', 'danger')
            return redirect(url_for('register'))
        
        if not any(c.isupper() for c in password):
            flash('❌ Password must contain at least one uppercase letter', 'danger')
            return redirect(url_for('register'))
        
        if not any(c.islower() for c in password):
            flash('❌ Password must contain at least one lowercase letter', 'danger')
            return redirect(url_for('register'))
        
        if not any(c.isdigit() for c in password):
            flash('❌ Password must contain at least one number', 'danger')
            return redirect(url_for('register'))
        
        if User.query.filter(User.username.ilike(username)).first():
            flash('❌ Username already exists', 'danger')
            return redirect(url_for('register'))
        
        if User.query.filter(User.email.ilike(email)).first():
            flash('❌ Email already registered', 'danger')
            return redirect(url_for('register'))
        
        hashed_password = generate_password_hash(
            password, 
            method='pbkdf2:sha256',
            salt_length=16
        )
        
        new_user = User(username=username, email=email, password=hashed_password)
        
        try:
            db.session.add(new_user)
            db.session.commit()
            flash('✅ Registration successful! Please login.', 'success')
            print(f"✅ New user registered: {username}")
            return redirect(url_for('login'))
        except Exception as e:
            db.session.rollback()
            flash('❌ Registration failed. Please try again.', 'danger')
            print(f"❌ Registration error: {e}")
    
    return render_template('register.html')

@app.route('/login', methods=['GET', 'POST'])
@limiter.limit("10 per minute")
def login():
    if request.method == 'POST':
        username = request.form.get('username', '').strip()
        password = request.form.get('password', '')
        
        if not username or not password:
            flash('⚠️ Please enter both username and password', 'danger')
            return redirect(url_for('login'))
        
        user = User.query.filter(User.username.ilike(username)).first()
        
        if user and check_password_hash(user.password, password):
            login_user(user)
            flash(f'🎉 Welcome back, {user.username}!', 'success')
            print(f"✅ User logged in: {user.username}")
            return redirect(url_for('home'))
        else:
            flash('❌ Invalid username or password', 'danger')
            print(f"❌ Failed login attempt for: {username}")
    
    return render_template('login.html')

@app.route('/logout')
@login_required
def logout():
    username = current_user.username
    logout_user()
    flash(f'👋 Goodbye, {username}! See you soon.', 'info')
    return redirect(url_for('login'))

@app.route('/')
def home():
    return render_template('index.html')

@app.route('/dashboard')
@login_required
def dashboard():
    user_uploads = Upload.query.filter_by(user_id=current_user.id).order_by(Upload.uploaded_at.desc()).all()
    user_quizzes = QuizResult.query.filter_by(user_id=current_user.id).order_by(QuizResult.completed_at.desc()).limit(5).all()
    
    return render_template('dashboard.html', uploads=user_uploads, quizzes=user_quizzes)

@app.route('/about')
def about():
    return render_template('about.html')

@app.route('/help')
def help_page():
    return render_template('help.html')

@app.route('/summarize', methods=['POST'])
@login_required
@limiter.limit("20 per hour")
def summarize():
    text = ""
    file = request.files.get('file')
    filename_to_save = 'Text Input'
    
    form_text = request.form.get('text', '').strip()
    if form_text:
        text = form_text[:10000]
        print(f"📝 Text from form: {len(text)} chars")

    if file and file.filename:
        filename = secure_filename(file.filename)
        filename_to_save = filename
        
        if not allowed_file(filename):
            flash("❌ Invalid file type. Allowed: PDF, DOCX, PPTX, TXT, PNG, JPG, JPEG", 'danger')
            return redirect(url_for('home'))
        
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        
        try:
            file.save(filepath)
            ext = filename.rsplit('.', 1)[1].lower()
            print(f"📁 File uploaded: {filename} ({ext})")
            
            # ✅ معالجة الصور
            if ext in ['png', 'jpg', 'jpeg', 'bmp', 'tiff', 'gif', 'webp']:
                from utils.ocr import extract_text_from_image
                file_text = extract_text_from_image(filepath)
                if not file_text or len(file_text.strip()) < 10:
                    flash("⚠️ Could not extract text from image. Make sure Tesseract OCR is installed, or paste the text manually.", 'warning')
                    return redirect(url_for('home'))
            else:
                file_text = extract_text_from_file(filepath, ext)
            
            if file_text:
                text = file_text
            else:
                flash("⚠️ Could not extract text from file", 'warning')
                return redirect(url_for('home'))
                
        except Exception as e:
            print(f"❌ Error processing file: {str(e)}")
            import traceback
            traceback.print_exc()
            flash(f"❌ Error processing file: {str(e)}", 'danger')
            return redirect(url_for('home'))
        finally:
            try:
                if os.path.exists(filepath):
                    os.remove(filepath)
                    print(f"🗑️ Cleaned up: {filename}")
            except:
                pass
    
    if not text or len(text.strip()) < 20:
        print("❌ No text or text too short")
        flash("⚠️ Please provide text or upload a file with sufficient content (at least 20 characters)", 'warning')
        return redirect(url_for('home'))

    if len(text) > 10000:
        text = text[:10000]
        flash("ℹ️ Text was truncated to 10,000 characters for optimal performance", 'info')

    print(f"📊 Final text length: {len(text)} chars")

    try:
        summary_style = request.form.get('summary_style', 'paragraphs')
        summary_length = request.form.get('summary_length', 'medium')
        chapter_id = request.form.get('chapter_id')
        
        print(f"🎨 Summary style: {summary_style}, length: {summary_length}")
        if chapter_id:
            print(f"📂 Saving to chapter_id: {chapter_id}")
        
        cleaned_text = preprocess_text(text)
        print(f"✅ Cleaned text: {len(cleaned_text)} chars")
        
        if len(cleaned_text) < 50:
            flash("⚠️ Text is too short after processing. Please provide more content.", 'warning')
            return redirect(url_for('home'))
        
        summary = summarize_text_with_style(cleaned_text, style=summary_style, length=summary_length)
        print(f"✅ Summary ({summary_style}, {summary_length}): {len(summary)} chars")
        
        if not summary or len(summary.strip()) < 10:
            flash("⚠️ Could not generate a meaningful summary. Please provide more content.", 'warning')
            return redirect(url_for('home'))
        
        quiz = generate_quiz(summary)
        print(f"✅ Quiz: {len(quiz)} questions")
        
        if not quiz or len(quiz) == 0:
            flash("⚠️ Could not generate quiz questions. Please try different content.", 'warning')
            return redirect(url_for('home'))
        
        upload = Upload(
            filename=filename_to_save,
            summary=summary,
            user_id=current_user.id,
            chapter_id=int(chapter_id) if chapter_id else None
        )
        db.session.add(upload)
        db.session.commit()
        
        session['summary'] = summary
        session['quiz'] = quiz
        session['score'] = 0

        def create_mindmap_async():
            try:
                words = summary.split()
                important_words = [
                    word.strip('.,;:!?"\'•1234567890') 
                    for word in words 
                    if len(word) > 4 and word.isalpha()
                ]
                
                if len(important_words) < 6:
                    important_words = [
                        word.strip('.,;:!?"\'•1234567890') 
                        for word in words 
                        if word.isalpha()
                    ][:6]
                else:
                    important_words = important_words[:6]
                
                create_mindmap(important_words)
                print(f"✅ Mindmap created asynchronously")
            except Exception as e:
                print(f"⚠️ Mindmap creation failed: {e}")
        
        executor.submit(create_mindmap_async)
        
        print("🎉 SUCCESS! Redirecting to results...")
        flash('✨ Summary and quiz generated successfully!', 'success')
        return render_template('result.html', summary=summary, quiz=quiz)
        
    except Exception as e:
        print(f"❌ Processing error: {str(e)}")
        import traceback
        traceback.print_exc()
        flash(f"❌ Error processing content: {str(e)}", 'danger')
        return redirect(url_for('home'))
# ================= Subjects & Chapters Routes =================
@app.route('/subjects')
@login_required
def subjects():
    user_subjects = Subject.query.filter_by(user_id=current_user.id).order_by(Subject.created_at.desc()).all()
    return render_template('subjects.html', subjects=user_subjects)

@app.route('/subject/create', methods=['GET', 'POST'])
@login_required
def create_subject():
    if request.method == 'POST':
        name = request.form.get('name', '').strip()
        description = request.form.get('description', '').strip()
        icon = request.form.get('icon', '📘')
        color = request.form.get('color', '#5DBAA4')
        
        if not name:
            flash('❌ Subject name is required', 'danger')
            return redirect(url_for('create_subject'))
        
        subject = Subject(
            name=name,
            description=description,
            icon=icon,
            color=color,
            user_id=current_user.id
        )
        
        try:
            db.session.add(subject)
            db.session.commit()
            flash(f'✅ Subject "{name}" created successfully!', 'success')
            return redirect(url_for('subjects'))
        except Exception as e:
            db.session.rollback()
            flash('❌ Error creating subject', 'danger')
            print(f"❌ Subject creation error: {e}")
    
    return render_template('create_subject.html')

@app.route('/subject/<int:subject_id>')
@login_required
def view_subject(subject_id):
    subject = Subject.query.get_or_404(subject_id)
    
    if subject.user_id != current_user.id:
        flash('❌ You do not have access to this subject', 'danger')
        return redirect(url_for('subjects'))
    
    chapters = Chapter.query.filter_by(subject_id=subject_id).order_by(Chapter.order).all()
    
    chapter_stats = {}
    for chapter in chapters:
        chapter_stats[chapter.id] = Upload.query.filter_by(chapter_id=chapter.id).count()
    
    return render_template('view_subject.html', subject=subject, chapters=chapters, chapter_stats=chapter_stats)

@app.route('/subject/<int:subject_id>/chapter/create', methods=['GET', 'POST'])
@login_required
def create_chapter(subject_id):
    subject = Subject.query.get_or_404(subject_id)
    
    if subject.user_id != current_user.id:
        flash('❌ You do not have access to this subject', 'danger')
        return redirect(url_for('subjects'))
    
    if request.method == 'POST':
        name = request.form.get('name', '').strip()
        description = request.form.get('description', '').strip()
        order = request.form.get('order', 1)
        
        if not name:
            flash('❌ Chapter name is required', 'danger')
            return redirect(url_for('create_chapter', subject_id=subject_id))
        
        chapter = Chapter(
            name=name,
            description=description,
            order=int(order),
            subject_id=subject_id
        )
        
        try:
            db.session.add(chapter)
            db.session.commit()
            flash(f'✅ Chapter "{name}" created successfully!', 'success')
            return redirect(url_for('view_subject', subject_id=subject_id))
        except Exception as e:
            db.session.rollback()
            flash('❌ Error creating chapter', 'danger')
            print(f"❌ Chapter creation error: {e}")
    
    last_chapter = Chapter.query.filter_by(subject_id=subject_id).order_by(Chapter.order.desc()).first()
    next_order = (last_chapter.order + 1) if last_chapter else 1
    
    return render_template('create_chapter.html', subject=subject, next_order=next_order)

@app.route('/chapter/<int:chapter_id>')
@login_required
def view_chapter(chapter_id):
    chapter = Chapter.query.get_or_404(chapter_id)
    subject = chapter.subject
    
    if subject.user_id != current_user.id:
        flash('❌ You do not have access to this chapter', 'danger')
        return redirect(url_for('subjects'))
    
    uploads = Upload.query.filter_by(chapter_id=chapter_id).order_by(Upload.uploaded_at.desc()).all()
    
    return render_template('view_chapter.html', chapter=chapter, subject=subject, uploads=uploads)

@app.route('/subject/<int:subject_id>/delete', methods=['POST'])
@login_required
def delete_subject(subject_id):
    subject = Subject.query.get_or_404(subject_id)
    
    if subject.user_id != current_user.id:
        flash('❌ You do not have access to this subject', 'danger')
        return redirect(url_for('subjects'))
    
    name = subject.name
    
    try:
        db.session.delete(subject)
        db.session.commit()
        flash(f'✅ Subject "{name}" deleted successfully', 'success')
    except Exception as e:
        db.session.rollback()
        flash('❌ Error deleting subject', 'danger')
        print(f"❌ Subject deletion error: {e}")
    
    return redirect(url_for('subjects'))

@app.route('/chapter/<int:chapter_id>/delete', methods=['POST'])
@login_required
def delete_chapter(chapter_id):
    chapter = Chapter.query.get_or_404(chapter_id)
    subject = chapter.subject
    
    if subject.user_id != current_user.id:
        flash('❌ You do not have access to this chapter', 'danger')
        return redirect(url_for('subjects'))
    
    name = chapter.name
    subject_id = chapter.subject_id
    
    try:
        db.session.delete(chapter)
        db.session.commit()
        flash(f'✅ Chapter "{name}" deleted successfully', 'success')
    except Exception as e:
        db.session.rollback()
        flash('❌ Error deleting chapter', 'danger')
        print(f"❌ Chapter deletion error: {e}")
    
    return redirect(url_for('view_subject', subject_id=subject_id))

# ================= Quiz Routes =================
@app.route('/quiz', methods=['GET', 'POST'])
@login_required
def quiz_page():
    if 'quiz' not in session:
        flash('⚠️ Please generate a summary first', 'warning')
        return redirect(url_for('home'))
    
    quiz = session['quiz']
    
    if request.method == 'POST':
        score = 0
        
        for i, q in enumerate(quiz):
            q_type = q.get('type', 'fill_blank')
            
            if q_type == 'mcq':
                user_answer = request.form.get(f'answer_{i}', '').strip()
                if user_answer.lower() == q['a'].lower():
                    score += 10
            
            elif q_type == 'true_false':
                user_answer = request.form.get(f'answer_{i}', '').strip().lower()
                if user_answer == q['a'].lower():
                    score += 10
            
            elif q_type == 'fill_blank':
                user_answer = request.form.get(f'answer_{i}', '').strip()
                if user_answer.lower() == q['a'].lower():
                    score += 10
            
            elif q_type == 'matching':
                correct_matches = 0
                for j in range(len(q.get('list_a', []))):
                    user_match = request.form.get(f'answer_{i}_{j}', '').strip()
                    correct_answer = q['a'].get(j, '')
                    if user_match == correct_answer:
                        correct_matches += 1
                
                if len(q.get('list_a', [])) > 0:
                    match_score = int((correct_matches / len(q['list_a'])) * 10)
                    score += match_score
        
        try:
            quiz_result = QuizResult(
                score=score,
                total_questions=len(quiz),
                user_id=current_user.id
            )
            db.session.add(quiz_result)
            
            current_user.score += score
            current_user.level = (current_user.score // 100) + 1
            db.session.commit()
            
            session['score'] = score
            flash(f'🎉 Quiz completed! You scored {score} points!', 'success')
            return redirect(url_for('gamification'))
        except Exception as e:
            db.session.rollback()
            flash('❌ Error saving quiz results', 'danger')
            print(f"❌ Quiz save error: {e}")
    
    return render_template('quiz.html', quiz=quiz)

# ================= Flashcards Routes =================
@app.route('/flashcards')
@login_required
def flashcards():
    if 'quiz' not in session:
        flash('⚠️ Please generate a summary first to create flashcards', 'warning')
        return redirect(url_for('home'))
    
    quiz = session['quiz']
    return render_template('flashcards.html', quiz=quiz)

@app.route('/submit-flashcards', methods=['POST'])
@login_required
def submit_flashcards():
    if 'quiz' not in session:
        return redirect(url_for('home'))
    
    known_count = int(request.form.get('known_count', 0))
    total_cards = len(session['quiz'])
    score = known_count * 10
    
    try:
        quiz_result = QuizResult(
            score=score,
            total_questions=total_cards,
            user_id=current_user.id
        )
        db.session.add(quiz_result)
        
        current_user.score += score
        current_user.level = (current_user.score // 100) + 1
        db.session.commit()
        
        flash(f'🎴 Flashcards completed! You scored {score} points!', 'success')
    except Exception as e:
        db.session.rollback()
        flash('❌ Error saving flashcard results', 'danger')
        print(f"❌ Flashcard save error: {e}")
    
    return redirect(url_for('gamification'))

# ================= Gamification Routes =================
@app.route('/gamification')
@login_required
def gamification():
    quiz = session.get('quiz', [])
    score = current_user.score
    return render_template('gamification.html', quiz=quiz, score=score)

@app.route('/leaderboard')
def leaderboard():
    top_users = User.query.order_by(User.score.desc()).limit(10).all()
    return render_template('leaderboard.html', users=top_users)

@app.route('/shared-library')
@login_required
def shared_library():
    shared_uploads = Upload.query.filter_by(is_shared=True).order_by(Upload.uploaded_at.desc()).all()
    return render_template('shared_library.html', uploads=shared_uploads)

@app.route('/share-upload/<int:upload_id>')
@login_required
def share_upload(upload_id):
    upload = Upload.query.get_or_404(upload_id)
    if upload.user_id == current_user.id:
        upload.is_shared = not upload.is_shared
        db.session.commit()
        status = 'shared' if upload.is_shared else 'unshared'
        flash(f'✅ Upload {status} successfully', 'success')
    else:
        flash('❌ You can only share your own uploads', 'danger')
    return redirect(url_for('dashboard'))

@app.route('/daily-challenge')
@login_required
def daily_challenge():
    today = datetime.utcnow().date()
    challenge = DailyChallenge.query.filter_by(date=today, active=True).first()
    
    if not challenge:
        challenge = DailyChallenge(
            title="Complete 3 Quizzes",
            description="Complete at least 3 quizzes today to earn bonus points!",
            points=50,
            date=today
        )
        db.session.add(challenge)
        db.session.commit()
    
    return render_template('daily_challenge.html', challenge=challenge)

# ================= Quiz Battles Routes =================
@app.route('/quiz-battles')
@login_required
def quiz_battles():
    battles = QuizBattle.query.filter_by(status='active').all()
    has_quiz = 'quiz' in session and len(session.get('quiz', [])) > 0
    return render_template('quiz_battles.html', battles=battles, has_quiz=has_quiz)

@app.route('/create-battle', methods=['POST'])
@login_required
def create_battle():
    if 'quiz' not in session or len(session.get('quiz', [])) == 0:
        flash('⚠️ Please generate a quiz first before creating a battle!', 'warning')
        return redirect(url_for('home'))
    
    title = request.form.get('title', 'Quick Battle')
    
    try:
        battle = QuizBattle(title=title)
        db.session.add(battle)
        db.session.commit()
        
        participant = BattleParticipant(battle_id=battle.id, user_id=current_user.id)
        db.session.add(participant)
        db.session.commit()
        
        flash('⚔️ Battle created successfully!', 'success')
        return redirect(url_for('battle_room', battle_id=battle.id))
    except Exception as e:
        db.session.rollback()
        flash('❌ Error creating battle', 'danger')
        print(f"❌ Battle creation error: {e}")
        return redirect(url_for('quiz_battles'))

@app.route('/battle/<int:battle_id>')
@login_required
def battle_room(battle_id):
    battle = QuizBattle.query.get_or_404(battle_id)
    quiz = session.get('quiz', [])
    
    if not quiz or len(quiz) == 0:
        flash('⚠️ Please generate a quiz first before joining battles!', 'warning')
        return redirect(url_for('home'))
    
    participant = BattleParticipant.query.filter_by(
        battle_id=battle_id, 
        user_id=current_user.id
    ).first()
    
    if not participant:
        participant = BattleParticipant(battle_id=battle_id, user_id=current_user.id)
        db.session.add(participant)
        db.session.commit()
    
    return render_template('battle_room.html', battle=battle, quiz=quiz)

@app.route('/submit-battle/<int:battle_id>', methods=['POST'])
@login_required
def submit_battle(battle_id):
    participant = BattleParticipant.query.filter_by(
        battle_id=battle_id,
        user_id=current_user.id
    ).first_or_404()
    
    quiz = session.get('quiz', [])
    
    if not quiz:
        flash('⚠️ Quiz not found in session', 'danger')
        return redirect(url_for('quiz_battles'))
    
    score = 0
    
    for i, q in enumerate(quiz):
        user_answer = request.form.get(f'answer_{i}', '').strip().lower()
        if user_answer == q['a'].lower():
            score += 10
    
    participant.score = score
    participant.completed = True
    
    current_user.score += score
    current_user.level = (current_user.score // 100) + 1
    
    db.session.commit()
    
    battle = QuizBattle.query.get(battle_id)
    all_completed = all(p.completed for p in battle.participants)
    
    if all_completed:
        winner = max(battle.participants, key=lambda p: p.score)
        battle.winner_id = winner.user_id
        battle.status = 'completed'
        db.session.commit()
    
    flash(f'⚔️ Battle submitted! Your score: {score}', 'success')
    return redirect(url_for('quiz_battles'))

# ================= Puzzle Mode Routes =================
@app.route('/puzzle-mode')
@login_required
def puzzle_mode():
    summary = session.get('summary', 'Learning is fun. Knowledge is power. Education transforms lives.')
    
    import random
    words = summary.split()
    scrambled = words.copy()
    random.shuffle(scrambled)
    
    puzzle = PuzzleGame(
        user_id=current_user.id,
        content=summary,
        scrambled_words=' '.join(scrambled)
    )
    db.session.add(puzzle)
    db.session.commit()
    
    return render_template('puzzle_mode.html', 
                         original=summary, 
                         scrambled=scrambled,
                         puzzle_id=puzzle.id)

@app.route('/submit-puzzle/<int:puzzle_id>', methods=['POST'])
@login_required
def submit_puzzle(puzzle_id):
    puzzle = PuzzleGame.query.get_or_404(puzzle_id)
    user_answer = request.form.get('answer', '').strip()
    
    original_text = re.sub(r'[^\w\s]', '', puzzle.content.lower())
    user_text = re.sub(r'[^\w\s]', '', user_answer.lower())
    
    original_words = original_text.split()
    user_words = user_text.split()
    
    correct = 0
    original_copy = original_words.copy()
    
    for word in user_words:
        if word in original_copy:
            correct += 1
            original_copy.remove(word)
    
    if len(original_words) > 0:
        percentage = (correct / len(original_words)) * 100
        score = int(percentage)
    else:
        score = 0
    
    puzzle.score = score
    puzzle.completed = True
    
    current_user.score += score
    current_user.level = (current_user.score // 100) + 1
    
    try:
        db.session.commit()
        flash(f'🧩 Puzzle completed! Score: {score}% ({correct}/{len(original_words)} words correct)', 'success')
    except Exception as e:
        db.session.rollback()
        flash('❌ Error saving puzzle results', 'danger')
        print(f"❌ Puzzle save error: {e}")
    
    return redirect(url_for('gamification'))

# ================= Grading Center Routes =================
@app.route('/grading-center')
@login_required
def grading_center():
    results = QuizResult.query.filter_by(user_id=current_user.id).order_by(
        QuizResult.completed_at.desc()
    ).all()
    
    total_quizzes = len(results)
    avg_score = sum(r.score for r in results) / total_quizzes if total_quizzes > 0 else 0
    
    return render_template('grading_center.html', 
                         results=results,
                         total_quizzes=total_quizzes,
                         avg_score=avg_score)

# ================= Review System Routes =================
@app.route('/review/<int:upload_id>', methods=['GET', 'POST'])
@login_required
def review_upload(upload_id):
    upload = Upload.query.get_or_404(upload_id)
    
    if request.method == 'POST':
        rating = int(request.form.get('rating', 3))
        comment = request.form.get('comment', '').strip()
        
        review = Review(
            upload_id=upload_id,
            user_id=current_user.id,
            rating=rating,
            comment=comment
        )
        db.session.add(review)
        db.session.commit()
        
        flash('⭐ Review submitted successfully!', 'success')
        return redirect(url_for('review_upload', upload_id=upload_id))
    
    reviews = Review.query.filter_by(upload_id=upload_id).all()
    avg_rating = sum(r.rating for r in reviews) / len(reviews) if reviews else 0
    
    return render_template('review_upload.html', 
                         upload=upload, 
                         reviews=reviews,
                         avg_rating=avg_rating)

# ================= Search Routes =================
@app.route('/search')
@login_required
def search():
    query = request.args.get('q', '').strip()
    
    if not query:
        return render_template('search.html', results=[], query='')
    
    results = Upload.query.filter(
        db.or_(
            Upload.filename.contains(query),
            Upload.summary.contains(query)
        )
    ).filter_by(is_shared=True).all()
    
    return render_template('search.html', results=results, query=query)

# ================= Export Routes =================
@app.route('/export/<int:upload_id>/<format>')
@login_required
def export_upload(upload_id, format):
    upload = Upload.query.get_or_404(upload_id)
    
    if format == 'txt':
        return Response(
            upload.summary,
            mimetype='text/plain',
            headers={'Content-Disposition': f'attachment;filename={upload.filename}.txt'}
        )
    
    flash('⚠️ Export format not supported yet', 'warning')
    return redirect(url_for('dashboard'))

# ================= Analytics Routes =================
@app.route('/analytics')
@login_required
def analytics():
    quiz_results = QuizResult.query.filter_by(user_id=current_user.id).all()
    uploads = Upload.query.filter_by(user_id=current_user.id).all()
    
    stats = calculate_user_stats(current_user, quiz_results, uploads)
    performance_data = analyze_quiz_performance(quiz_results)
    weekly_report = generate_performance_report(current_user, quiz_results)
    
    return render_template('analytics.html', 
                         stats=stats,
                         performance=performance_data,
                         weekly_report=weekly_report)

# ================= Visualization Routes =================
@app.route('/visualize/<int:upload_id>')
@login_required
def visualize_upload(upload_id):
    upload = Upload.query.get_or_404(upload_id)
    
    keywords = extract_keywords(upload.summary)
    keyword_freq = {kw: upload.summary.lower().count(kw) for kw in keywords}
    
    chart_path = create_bar_chart(
        list(keyword_freq.keys()),
        list(keyword_freq.values()),
        f"Keyword Frequency: {upload.filename}"
    )
    
    return render_template('visualize.html', 
                         upload=upload,
                         chart=chart_path,
                         keywords=keywords)

# ================= Admin Routes =================
@app.route('/admin/login', methods=['GET', 'POST'])
@limiter.limit("3 per minute")
def admin_login():
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        
        if username == 'admin' and password == 'vortex2026':
            session['is_admin'] = True
            flash('👨‍💼 Admin login successful!', 'success')
            return redirect(url_for('admin_panel'))
        else:
            flash('❌ Invalid admin credentials', 'danger')
    
    return render_template('admin_login.html')

@app.route('/admin/panel')
def admin_panel():
    if not session.get('is_admin'):
        flash('❌ Admin access required', 'danger')
        return redirect(url_for('admin_login'))
    
    total_users = User.query.count()
    total_uploads = Upload.query.count()
    total_quizzes = QuizResult.query.count()
    total_battles = QuizBattle.query.count()
    total_subjects = Subject.query.count()
    total_chapters = Chapter.query.count()
    
    recent_users = User.query.order_by(User.created_at.desc()).limit(10).all()
    recent_uploads = Upload.query.order_by(Upload.uploaded_at.desc()).limit(10).all()
    
    return render_template('admin_panel.html',
                         total_users=total_users,
                         total_uploads=total_uploads,
                         total_quizzes=total_quizzes,
                         total_battles=total_battles,
                         total_subjects=total_subjects,
                         total_chapters=total_chapters,
                         recent_users=recent_users,
                         recent_uploads=recent_uploads)

@app.route('/admin/logout')
def admin_logout():
    session.pop('is_admin', None)
    flash('👋 Admin logged out', 'info')
    return redirect(url_for('home'))

@app.route('/admin/users')
def admin_users():
    if not session.get('is_admin'):
        return redirect(url_for('admin_login'))
    
    users = User.query.order_by(User.score.desc()).all()
    return render_template('admin_users.html', users=users)

@app.route('/admin/delete-user/<int:user_id>')
def admin_delete_user(user_id):
    if not session.get('is_admin'):
        return redirect(url_for('admin_login'))
    
    user = User.query.get_or_404(user_id)
    username = user.username
    
    Upload.query.filter_by(user_id=user_id).delete()
    QuizResult.query.filter_by(user_id=user_id).delete()
    
    db.session.delete(user)
    db.session.commit()
    
    flash(f'✅ User {username} deleted successfully', 'success')
    return redirect(url_for('admin_users'))
# ================= Knowledge Base Routes =================
@app.route('/knowledge-base')
@login_required
def knowledge_base():
    categories = get_all_categories()
    category_topics = {}
    for cat in categories:
        category_topics[cat] = get_topics_by_category(cat)
    
    return render_template('knowledge_base.html', 
                         categories=categories,
                         category_topics=category_topics)

@app.route('/knowledge-base/topic/<topic_id>')
@login_required
def view_knowledge_topic(topic_id):
    topic = get_topic_content(topic_id)
    
    if not topic:
        flash('❌ Topic not found', 'danger')
        return redirect(url_for('knowledge_base'))
    
    return render_template('view_knowledge_topic.html', topic=topic, topic_id=topic_id)

@app.route('/knowledge-base/use/<topic_id>', methods=['POST'])
@login_required
@limiter.limit("20 per hour")
def use_knowledge_topic(topic_id):
    topic = get_topic_content(topic_id)
    
    if not topic:
        flash('❌ Topic not found', 'danger')
        return redirect(url_for('knowledge_base'))
    
    try:
        text = topic['content']
        
        summary_style = request.form.get('summary_style', 'paragraphs')
        summary_length = request.form.get('summary_length', 'medium')
        chapter_id = request.form.get('chapter_id')
        
        cleaned_text = preprocess_text(text)
        summary = summarize_text_with_style(cleaned_text, style=summary_style, length=summary_length)
        quiz = generate_quiz(summary)
        
        upload = Upload(
            filename=f"📚 {topic['title']}",
            summary=summary,
            user_id=current_user.id,
            chapter_id=int(chapter_id) if chapter_id else None
        )
        db.session.add(upload)
        db.session.commit()
        
        session['summary'] = summary
        session['quiz'] = quiz
        session['score'] = 0
        
        flash('✨ Knowledge topic processed successfully!', 'success')
        return render_template('result.html', summary=summary, quiz=quiz)
        
    except Exception as e:
        flash(f'❌ Error processing topic: {str(e)}', 'danger')
        return redirect(url_for('knowledge_base'))

@app.route('/knowledge-base/search')
@login_required
def search_knowledge():
    query = request.args.get('q', '').strip()
    
    if not query:
        return redirect(url_for('knowledge_base'))
    
    results = search_topics(query)
    
    return render_template('search_knowledge.html', results=results, query=query)
# ================= Error Handlers =================
@app.errorhandler(404)
def not_found_error(error):
    return render_template('404.html'), 404

@app.errorhandler(500)
def internal_error(error):
    db.session.rollback()
    return render_template('500.html'), 500

@app.errorhandler(429)
def ratelimit_handler(e):
    flash('⚠️ Too many requests! Please wait a moment and try again.', 'warning')
    return render_template('429.html', error=e), 429

# ================= Run Server =================
if __name__ == '__main__':
    print("=" * 60)
    print("🚀 VORTEX - AI-Powered Learning Platform")
    print("=" * 60)
    print("✨ Pastel Illustration Theme Enabled")
    print("⚡ Performance Optimizations Active")
    print("🎨 Summary Customization Enabled")
    print("📂 Subjects & Chapters System Active")
    print("🔒 Security Features Active:")
    print("   • Password Hashing (PBKDF2-SHA256 + Salt)")
    print("   • Rate Limiting (5/min Register, 10/min Login)")
    print("   • Enhanced Validation (Email, Username, Password)")
    print("   • Case-Insensitive User Lookup")
    print("   • Generic Error Messages")
    print("📚 All Features Active:")
    print("   • Smart Summarization (5 Styles, 3 Lengths)")
    print("   • Fast File Processing (32MB, Async)")
    print("   • Subjects & Chapters Organization")
    print("   • Detailed Mind Maps (25 Nodes)")
    print("   • Interactive Quizzes")
    print("   • Flashcards System")
    print("   • Gamification")
    print("   • Quiz Battles")
    print("   • Puzzle Mode")
    print("   • Analytics Dashboard")
    print("   • Admin Panel (Rate Limited)")
    print("=" * 60)
    print("🌐 Server running on: http://localhost:5000")
    print("💜 Developed by Vortex Team - Level 10")
    print("=" * 60)
    app.run(host="0.0.0.0", port=5000, debug=True)